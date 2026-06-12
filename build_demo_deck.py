"""
Build a 28-slide PowerPoint of Cesar's Amira demo screenshots.

Each slide:
- White background
- Step label header at top (bold, dark charcoal, ~24pt)
- Screenshot centered below, scaled to fit while preserving aspect ratio

16:9 widescreen format (13.33" x 7.5").
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

SHOTS_DIR = r"C:\Users\farza\Desktop\new screenshots\demo-screenshots"
OUT_PATH = r"C:\Users\farza\Desktop\Amira_Demo_Walkthrough_2026-04-28.pptx"

# (filename_stem, step_number, step_title) — order matters, this is presentation order
SLIDES = [
    ("step-01-home-portal",              "1",   "Home portal"),
    ("step-02a-spec-new",                "2a",  "Spec Agent entry — natural-language intent"),
    ("step-02b-projects-new",            "2b",  "Structured project creation"),
    ("step-02c-import-from-repo",        "2c",  "Import from repository"),
    ("step-03-spec-finiq",               "3",   "Decision Point — alternatives with trade-offs"),
    ("step-04-version-history",          "4",   "Living specification, gaps, version history"),
    ("step-05-route-esignature-modal",   "5",   "Route specification for e-signature"),
    ("step-06a-approve-page",            "6a",  "Authorized approver view"),
    ("step-06b-signature-recorded",      "6b",  "Signature recorded"),
    ("step-07-canvas-preview",           "7",   "Open Canvas — three-panel workspace"),
    ("step-08a-enlarge-chart-clicked",   "8a",  "Build Agent — enlarge chart"),
    ("step-08b-working-capital-added",   "8b",  "Build Agent — add Working Capital KPI"),
    ("step-09a-resources-tab",           "9a",  "Resources tab — bound primitives"),
    ("step-09b-add-resource-drawer",     "9b",  "Add to App — Skills"),
    ("step-09c-add-companion-agent",     "9c",  "Add to App — Companion Agents"),
    ("step-09d-add-knowledge",           "9d",  "Add to App — Knowledge files"),
    ("step-10-compliance-matrix",        "10",  "Compliance Matrix"),
    ("step-11a-deploy-publish",          "11a", "Deploy — Publish details"),
    ("step-11b-deploy-compliance",       "11b", "Deploy — Compliance & Governance"),
    ("step-11c-deploy-environment",      "11c", "Deploy — Environment"),
    ("step-11d-deploy-approval",         "11d", "Deploy — Approval & e-signature"),
    ("step-11e-deploy-progress",         "11e", "Deploying…"),
    ("step-12a-skills",                  "12a", "Skills marketplace"),
    ("step-12b-skills-app-agents",       "12b", "Apps Become Agents"),
    ("step-12c-skills-bottom",           "12c", "Skills marketplace — full catalog"),
    ("step-13a-ask-amira-drawer",        "13a", "Ask Amira — FinIQ Agent (Q3 Petcare)"),
    ("step-13b-ask-amira-nestle",        "13b", "Ask Amira — comparison vs Nestlé"),
    ("step-14-project-finiq",            "14",  "Project FinIQ — lineage and audit"),
]

# Slide geometry (16:9 widescreen)
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Header band geometry
HEADER_H_IN = 0.85
HEADER_PAD_LEFT_IN = 0.6

# Image area (below the header)
IMG_AREA_TOP_IN = HEADER_H_IN + 0.1  # small gap below the header
IMG_AREA_LEFT_IN = 0.4
IMG_AREA_RIGHT_IN = 0.4
IMG_AREA_BOTTOM_IN = 0.3
IMG_AREA_W_IN = SLIDE_W_IN - IMG_AREA_LEFT_IN - IMG_AREA_RIGHT_IN
IMG_AREA_H_IN = SLIDE_H_IN - IMG_AREA_TOP_IN - IMG_AREA_BOTTOM_IN

# Colors — Charcoal Minimal palette (per skill guidance: charcoal header, white body)
CHARCOAL = RGBColor(0x36, 0x45, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x21, 0x21, 0x21)


def fit_dimensions(img_w_px, img_h_px, area_w_in, area_h_in):
    """Scale (img_w_px, img_h_px) to fit inside (area_w_in, area_h_in) preserving aspect ratio.
    Returns (width_in, height_in)."""
    aspect = img_w_px / img_h_px
    area_aspect = area_w_in / area_h_in
    if aspect > area_aspect:
        # Image is wider — fit to width
        return area_w_in, area_w_in / aspect
    else:
        # Image is taller — fit to height
        return area_h_in * aspect, area_h_in


def add_step_slide(prs, image_path, step_num, step_title):
    """Add one slide with header band + step label + centered screenshot."""
    blank_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank_layout)

    # 1. Header band — charcoal rectangle across the top
    header_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(SLIDE_W_IN), Inches(HEADER_H_IN),
    )
    header_shape.line.fill.background()  # no border
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = CHARCOAL
    # Make sure no shadow
    sp = header_shape.shadow
    sp.inherit = False

    # 2. Step label text on the header (e.g. "Step 7 — Open Canvas — three-panel workspace")
    label_box = slide.shapes.add_textbox(
        Inches(HEADER_PAD_LEFT_IN), Inches(0),
        Inches(SLIDE_W_IN - HEADER_PAD_LEFT_IN - 0.4), Inches(HEADER_H_IN),
    )
    tf = label_box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = 3  # MSO_ANCHOR.MIDDLE
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    # "Step 7" in slightly smaller bold + " — " + title
    run_step = p.add_run()
    run_step.text = f"Step {step_num}"
    run_step.font.name = "Calibri"
    run_step.font.size = Pt(22)
    run_step.font.bold = True
    run_step.font.color.rgb = WHITE

    run_dash = p.add_run()
    run_dash.text = "   "
    run_dash.font.name = "Calibri"
    run_dash.font.size = Pt(22)
    run_dash.font.color.rgb = WHITE

    run_title = p.add_run()
    run_title.text = step_title
    run_title.font.name = "Calibri"
    run_title.font.size = Pt(22)
    run_title.font.bold = False
    run_title.font.color.rgb = WHITE

    # 3. Image — open with PIL to read pixel dimensions, scale to fit, center
    with Image.open(image_path) as im:
        px_w, px_h = im.size

    fit_w, fit_h = fit_dimensions(px_w, px_h, IMG_AREA_W_IN, IMG_AREA_H_IN)
    left_in = IMG_AREA_LEFT_IN + (IMG_AREA_W_IN - fit_w) / 2
    top_in = IMG_AREA_TOP_IN + (IMG_AREA_H_IN - fit_h) / 2

    slide.shapes.add_picture(
        image_path,
        Inches(left_in), Inches(top_in),
        width=Inches(fit_w), height=Inches(fit_h),
    )


def main():
    prs = Presentation()
    # Set slide size to 16:9 widescreen
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    missing = []
    for stem, step_num, step_title in SLIDES:
        path = os.path.join(SHOTS_DIR, f"{stem}.png")
        if not os.path.exists(path):
            missing.append(stem)
            continue
        add_step_slide(prs, path, step_num, step_title)

    if missing:
        print(f"WARNING: missing files: {missing}")

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
